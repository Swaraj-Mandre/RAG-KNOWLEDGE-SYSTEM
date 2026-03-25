import os
import time
from pathlib import Path
from dotenv import load_dotenv

from langchain_community.document_loaders import TextLoader, PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document 

load_dotenv(dotenv_path=Path(__file__).resolve().parent.parent / ".env", override=True)

DOCUMENTS_PATH = Path(__file__).resolve().parent.parent / "data" / "documents"


# Image (.png/.jpeg)
def load_image(file_path):
    """Extract text from image using Gemini Vision - FREE same quota"""
    try:
        import base64
        from google import genai

        api_key = os.getenv("GOOGLE_API_key")
        client = genai.Client(api_key=api_key)

        with open(file_path, "rb") as f:
            image_data = base64.b64encode(f.read()).decode("utf-8")
        suffix = Path(file_path).suffix.lower()
        mime = "image/png" if suffix == ".png" else "image/jpeg"

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[
                {"inline_data": {"mime_type": mime, "data": image_data}},
                "Extract ALL text and describe all content visible in this image in detail."
            ]
        )
        text = response.text
        return [Document(page_content=text, metadata={"source": Path(file_path).name, "type":"image"})]
    except Exception as e:
        print(f"    Could not process image {Path(file_path).name}: {e}")
        return[]

# PPTX
def load_pptx(file_path):
    """Extract text from PowerPoint slides"""
    try: 
        from pptx import Presentation
        prs = Presentation(file_path)
        text_chunks = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                content = f"Slide {i+1}:\n" + "\n".join(slide_text)
                text_chunks.append(Document(
                    page_content=content,
                    metadata={"source": Path(file_path).name, "slide": i+1}
                ))
        return text_chunks
    except Exception as e:
        print(f"   Could not process .pptx {Path(file_path).name}: {e}")
        return[]
    
# DOCX
def load_docx(file_path):
    """Extract text from Word Documents"""
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(file_path)
        full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return [Document(
            page_content=full_text,
            metadata={"source": Path(file_path).name, "type": "docx"}
        )]
    except Exception as e:
        print(f"     Could not process DOCX {Path(file_path).name}: {e}")
        return[]
    
def load_all_documents():
    all_documents = []
    supported = [".txt", ".pdf", ".png", ".jpeg", ".jpg", ".pptx", ".docx"]
    files_found = [f for f in DOCUMENTS_PATH.iterdir() if f.suffix.lower() in supported]

    if not files_found:
        raise FileNotFoundError(
            "No Documents found\n"
            "Place files in data/documents / folder"
        )
    
    print(f"Found {len(files_found)} file(s):")
    for file in files_found:
        try:
            print(f"   Loading: {file.name}")
            ext = file.suffix.lower()
            docs = []

            if ext == ".pdf":
                loader = PyPDFLoader(str(file))
            elif ext == ".txt":
                loader = TextLoader(str(file), encoding="utf-8") #utf-8 , saves program from crashing for special character or emojis
                docs = loader.load() #yet we got address of document '.load()' plays important role to go to harddrive, opem the file, read the text , and bring back into RAM (memory) 
            elif ext in [".png", ".jpg", ".jpeg"]:
                docs = load_image(str(file))
            elif ext == ".pptx":
                docs = load_pptx(str(file))
            elif ext == ".docx":
                dpcs = load_docx(str(file))
            else:
                continue

            for doc in docs : 
                doc.metadata["source"] = file.name
            all_documents.extend(docs) 
            print(f"      {file.name} - {len(docs)} section(s) loaded")  #len() = "How many iteams do you contains?" 
        except Exception as e:
            print(f"      Skipping {file.name}: {e}")
    return all_documents

#Capacity check 
def check_capacity(chunks):
    total = len(chunks)
    print(f"\nCapacity Check:")
    print(f"   Total chunks: {total}")

    if total <= 80:
        print("   Safe — will process instantly")
        return False  # no delay needed
    elif total <= 500:
        print("   Large — adding delays between batches")
        return True   # delay needed
    else:
        print("   Too large! Increase chunk_size to reduce chunks")
        print(f"  Try chunk_size=2000 to get ~{total//2} chunks")
        exit(1)

#def is function , in which anything created (like 'loader' 'chunks' 'vector_db') is LOCAL. Once function finishes, Python cleans up memory and those variable disapper.
def ingest_documents():
    print("Loading documents...")

    # 1. Load document
    documents = load_all_documents()
    print(f"\n Total Pages Loaded: {len(documents)}")

    # 2. Split into chunks
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000, #bigger = fewer chunks  = fewer API calls
        # char = 150 words = 0.5 page
        # 1000 chunks * 0.5 page = ~500 pages MAX per day
        chunk_overlap=100 #It's like going 50 piece back in last chunk and include that into new chunk . By which no gap , it didnt miss content flow
    )
    chunks = splitter.split_documents(documents)
    print(f"Created {len(chunks)} chunks")

    needs_delay = check_capacity(chunks)

    # 3. Create embeddings in batches to avoid rate limit 
    print("Creating embeddings...")
    embeddings = GoogleGenerativeAIEmbeddings(
        model="models/gemini-embedding-001"
    )

    # 4. Store in FAISS
    if needs_delay:
        batch_size = 80
        all_docs = []
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i+batch_size]
            print(f"  Processing batch {i//batch_size + 1} ({len(batch)} chunks). . .")
            all_docs.extend(batch)
            if i + batch_size < len(chunks):
                print("   Waiting 60s for rate limit reset. . .")
                time.sleep(60)
        vector_db = FAISS.from_documents(all_docs, embeddings)
    else:
        vector_db = FAISS.from_documents(chunks, embeddings)


    # 5. Save locally
    vector_db.save_local("vectorstore")
    print("Vector database created and saved!")
    print(f"  Location: vectorstore")

if __name__ == "__main__":
    ingest_documents()